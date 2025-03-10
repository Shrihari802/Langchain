import os
import streamlit as st
import PyPDF2
import docx
import pptx
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
from langchain_core.prompts import ChatPromptTemplate
from langchain_groq import ChatGroq

# Set up Groq API key
os.environ["GROQ_API_KEY"] = "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"

# Load model
model = SentenceTransformer("all-MiniLM-L6-v2")
llm = ChatGroq(temperature=0, model_name="deepseek-r1-distill-qwen-32b")

def extract_text_from_pdf(pdf_file):
    reader = PyPDF2.PdfReader(pdf_file)
    text = "".join([page.extract_text() for page in reader.pages if page.extract_text()])
    return text

def extract_text_from_docx(docx_file):
    doc = docx.Document(docx_file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(pptx_file):
    presentation = pptx.Presentation(pptx_file)
    text = "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
    return text

def extract_text_from_txt(txt_file):
    return txt_file.read().decode("utf-8")

def get_most_relevant_text(query, corpus, corpus_embeddings):
    query_embedding = model.encode([query])
    similarities = cosine_similarity(query_embedding, corpus_embeddings)[0]
    best_idx = similarities.argmax()
    return corpus[best_idx]

def generate_answer(query, context):
    prompt = ChatPromptTemplate.from_messages([
        ("system", "You are an assistant that answers questions based on provided text."),
        ("user", f"Context: {context}\n\nQuestion: {query}")
    ])

    formatted_prompt = prompt.format()  # Convert prompt to a string
    raw_response = llm.invoke(formatted_prompt)  # Get the full response

    # Extract only the relevant answer (from "Overview" onwards)
    if "OlmOCR Overview" in raw_response:
        answer = raw_response.split("OlmOCR Overview")[1]  # Get everything after "OlmOCR Overview"
        return "OlmOCR Overview" + answer  # Add back the heading
    else:
        return raw_response  # If no "Overview", return everything


# Streamlit UI
st.set_page_config(page_title="Multi-Format Q&A with RAG", layout="wide")
st.title("📄 Multi-Format Q&A with Retrieval-Augmented Generation")

uploaded_file = st.file_uploader("Upload a document", type=["pdf", "docx", "pptx", "txt"])

if uploaded_file:
    with st.spinner("Processing document..."):
        if uploaded_file.type == "application/pdf":
            document_text = extract_text_from_pdf(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            document_text = extract_text_from_docx(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            document_text = extract_text_from_pptx(uploaded_file)
        elif uploaded_file.type == "text/plain":
            document_text = extract_text_from_txt(uploaded_file)
        else:
            st.error("Unsupported file format!")
            document_text = ""
        
        if document_text:
            sentences = document_text.split(". ")  # Basic sentence splitting
            sentence_embeddings = model.encode(sentences)
            st.success("Document processed successfully!")
        
    query = st.text_input("Ask a question about the document:")
    if query and document_text:
        with st.spinner("Searching for relevant information..."):
            relevant_text = get_most_relevant_text(query, sentences, sentence_embeddings)
            answer = generate_answer(query, relevant_text)
            st.write("### Answer:")
            st.info(generate_answer(query, relevant_text))  # Show only cleaned-up answer
            
